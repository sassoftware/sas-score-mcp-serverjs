"""
make-pptx.py  —  Build a PPTX from slide-02.png and slide-03.png
Run: python make-pptx.py
Output: sas-score-mcp-overview.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Slide dimensions: 16:9 widescreen (10 × 5.625 inches)
W = Inches(10)
H = Inches(5.625)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

SRC_DIR = os.path.dirname(os.path.abspath(__file__))
images  = ["slide-02.png", "slide-03.png"]

for img_name in images:
    img_path = os.path.join(SRC_DIR, img_name)
    slide = prs.slides.add_slide(blank_layout)
    # Place the PNG as a full-bleed image covering the entire slide
    slide.shapes.add_picture(img_path, Inches(0), Inches(0), W, H)

out = os.path.join(SRC_DIR, "sas-score-mcp-overview.pptx")
prs.save(out)
print(f"Saved: {out}")
